#!/usr/bin/env python3
"""
AI4S 科研选题方法论 — 方法驱动 vs 问题驱动
Speaker: 张恒  |  4 slides  |  16:9 widescreen
v3: Larger fonts (body ≥10pt), compact text, generous spacing, no overlaps.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── Constants ──────────────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

USTC_BLUE  = RGBColor(0x00, 0x3D, 0x7C)
ACCENT_RED = RGBColor(0xCC, 0x00, 0x00)
LIGHT_BLUE = RGBColor(0xE6, 0xF0, 0xF8)
DARK_BLUE  = RGBColor(0x00, 0x2B, 0x5A)
LIGHT_GRAY = RGBColor(0xE8, 0xE8, 0xE8)
MED_GRAY   = RGBColor(0xAA, 0xAA, 0xAA)
DARK_GRAY  = RGBColor(0x55, 0x55, 0x55)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x00, 0x00, 0x00)
ORANGE     = RGBColor(0xE6, 0x7E, 0x22)
GREEN      = RGBColor(0x27, 0xAE, 0x60)
PURPLE     = RGBColor(0x8E, 0x44, 0xAD)

FONT = 'Noto Sans CJK SC'
LM  = Inches(0.85)
RM  = Inches(0.85)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]


# ── Low-level helpers ──────────────────────────────────────

def _font(run, sz, color, bold=False):
    run.font.size = Pt(sz)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = FONT
    rPr = run._r.get_or_add_rPr()
    rPr.set(qn('a:altLang'), 'zh-CN')
    ea = rPr.makeelement(qn('a:ea'), {})
    ea.set('typeface', FONT)
    rPr.append(ea)


def _tb(slide, l, t, w, h, text, sz=14, color=BLACK, bold=False,
        align=PP_ALIGN.LEFT, ls=1.35):
    """Single-paragraph textbox."""
    box = slide.shapes.add_textbox(l, t, w, h)
    box.word_wrap = True
    tf = box.text_frame; tf.word_wrap = True; tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text; p.alignment = align
    p.space_after = Pt(0); p.space_before = Pt(0)
    p.line_spacing = Pt(int(sz * ls))
    _font(p.runs[0] if p.runs else p.add_run(), sz, color, bold)
    return box


def _mtb(slide, l, t, w, h, lines, sz=10, color=DARK_GRAY, ls=1.35):
    """Multi-paragraph textbox. lines: (text, bold, color_or_none)."""
    box = slide.shapes.add_textbox(l, t, w, h)
    box.word_wrap = True
    tf = box.text_frame; tf.word_wrap = True; tf.auto_size = None
    for i, item in enumerate(lines):
        if isinstance(item, str):        txt, bld, clr = item, False, color
        elif len(item) == 2:             txt, bld, clr = item[0], item[1], color
        else:                            txt, bld, clr = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt; p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(1); p.space_before = Pt(1)
        p.line_spacing = Pt(int(sz * ls))
        if p.runs: _font(p.runs[0], sz, clr, bld)
    return box


def _rect(slide, l, t, w, h, fill=WHITE, border=None, bw=Pt(1.5), radius=None):
    st = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(st, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border: s.line.color.rgb = border; s.line.width = bw
    else:      s.line.fill.background()
    return s


def _line(slide, l, t, w, color=LIGHT_GRAY, lw=Pt(1)):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, Pt(0))
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background(); s.height = lw
    return s


def _arrow(slide, l, t, w, h, color=USTC_BLUE):
    s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def _strip(slide, l, t, w, h, text, bg=LIGHT_BLUE, tc=USTC_BLUE, sz=11, bold=True):
    box = _rect(slide, l, t, w, h, fill=bg, radius=Inches(0.05))
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.15); tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.06); tf.margin_bottom = Inches(0.06)
    p = tf.paragraphs[0]; p.text = text; p.alignment = PP_ALIGN.LEFT
    if p.runs: _font(p.runs[0], sz, tc, bold)


def _table(slide, l, t, col_w, rows, font_sz=9.5):
    nr, nc = len(rows), len(col_w)
    tw = sum(col_w)
    rh = Inches(0.38)   # row height
    ts = slide.shapes.add_table(nr, nc, l, t, tw, rh * nr)
    tbl = ts.table
    for ci, w in enumerate(col_w): tbl.columns[ci].width = w
    for ri, row in enumerate(rows):
        for ci, ct in enumerate(row):
            cell = tbl.cell(ri, ci); cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.text = str(ct); p.alignment = PP_ALIGN.LEFT
            p.space_before = Pt(2); p.space_after = Pt(2)
            p.line_spacing = Pt(int(font_sz * 1.25))
            is_hdr = (ri == 0)
            bld = is_hdr or (ci == 0)
            tc2 = WHITE if is_hdr else BLACK
            if p.runs: _font(p.runs[0], font_sz, tc2, bld)
            tcPr = cell._tc.get_or_add_tcPr()
            sf = cell._tc.makeelement(qn('a:solidFill'), {})
            val = '003D7C' if is_hdr else ('F5F8FC' if ri % 2 == 1 else 'FFFFFF')
            srgb = sf.makeelement(qn('a:srgbClr'), {'val': val})
            sf.append(srgb); tcPr.append(sf)
            for bt in ['a:lnL','a:lnR','a:lnT','a:lnB']:
                ln = cell._tc.makeelement(qn(bt), {'w':'6350'})
                sf2 = ln.makeelement(qn('a:solidFill'), {})
                srgb2 = sf2.makeelement(qn('a:srgbClr'), {'val':'D0D0D0'})
                sf2.append(srgb2); ln.append(sf2); tcPr.append(ln)
    return ts


# ── Common chrome ──────────────────────────────────────────

def add_header(slide):
    _rect(slide, Inches(0), Inches(0), Inches(0.08), Inches(0.48), fill=USTC_BLUE)
    _line(slide, Inches(0), Inches(0.48), SLIDE_W, USTC_BLUE, Pt(2.5))
    hw = Inches(5.5)
    _tb(slide, SLIDE_W - hw - Inches(0.5), Inches(0.06), hw, Inches(0.2),
        '中国科学技术大学', 10, USTC_BLUE, True, PP_ALIGN.RIGHT)
    _tb(slide, SLIDE_W - hw - Inches(0.5), Inches(0.25), hw, Inches(0.16),
        'University of Science and Technology of China', 6.5, MED_GRAY, False, PP_ALIGN.RIGHT)
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, SLIDE_W - Inches(0.9),
                                  Inches(0.11), Inches(0.05), Inches(0.05))
    dot.fill.solid(); dot.fill.fore_color.rgb = USTC_BLUE; dot.line.fill.background()


def add_title(slide, text):
    _tb(slide, LM, Inches(0.68), Inches(11), Inches(0.38), text, 21, USTC_BLUE, True)
    _line(slide, LM, Inches(1.08), Inches(3.2), USTC_BLUE, Pt(3))


def add_assertion(slide, text, y=Inches(1.28)):
    _arrow(slide, LM, y + Inches(0.03), Inches(0.25), Inches(0.18), USTC_BLUE)
    _tb(slide, LM + Inches(0.4), y, Inches(11.5), Inches(0.35), text, 12, DARK_BLUE, True)


def add_footer(slide, n):
    _line(slide, LM, Inches(6.9), SLIDE_W - LM - RM, LIGHT_GRAY, Pt(1))
    _tb(slide, SLIDE_W - RM - Inches(0.7), Inches(6.97), Inches(0.7), Inches(0.22),
        str(n), 8, MED_GRAY, False, PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ═══════════════════════════════════════════════════════════

def build_s1():
    s = prs.slides.add_slide(blank_layout)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE
    add_header(s)
    _rect(s, Inches(0), Inches(0.48), Inches(0.08), Inches(6.4), fill=USTC_BLUE)

    _tb(s, Inches(1.5), Inches(1.5), Inches(10.3), Inches(1.0),
        'AI4S 领域的科研选题', 46, USTC_BLUE, True, PP_ALIGN.CENTER)
    _line(s, Inches(3.8), Inches(2.55), Inches(5.7), USTC_BLUE, Pt(4))

    _tb(s, Inches(1.5), Inches(2.95), Inches(10.3), Inches(0.5),
        '方法驱动  vs  问题驱动', 24, DARK_BLUE, True, PP_ALIGN.CENTER)
    _tb(s, Inches(1.5), Inches(3.45), Inches(10.3), Inches(0.45),
        '两条路径的选择、差异与融合', 18, DARK_GRAY, False, PP_ALIGN.CENTER)

    _line(s, Inches(5.0), Inches(4.15), Inches(3.3), MED_GRAY, Pt(1))

    _tb(s, Inches(1.5), Inches(4.35), Inches(10.3), Inches(0.4),
        '张  恒', 22, USTC_BLUE, True, PP_ALIGN.CENTER)
    _tb(s, Inches(1.5), Inches(4.75), Inches(10.3), Inches(0.3),
        'Workshop: How to Conduct Research — 选题：构建领域视野', 12, DARK_GRAY, False,
        PP_ALIGN.CENTER)

    # Bottom tags
    tags  = ['路径选择', 'Gap 识别', '真伪判断', '决策框架']
    tcols = [USTC_BLUE, ORANGE, GREEN, PURPLE]
    tw = Inches(2.3); gap = Inches(0.3)
    total = len(tags) * tw + (len(tags) - 1) * gap
    sx = (SLIDE_W - total) / 2; ty = Inches(5.5)
    for i, (tag, tc) in enumerate(zip(tags, tcols)):
        b = _rect(s, sx + i * (tw + gap), ty, tw, Inches(0.4), fill=tc, radius=Inches(0.05))
        tf = b.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = tag; p.alignment = PP_ALIGN.CENTER
        if p.runs: _font(p.runs[0], 12, WHITE, True)

    add_footer(s, 1)


# ═══════════════════════════════════════════════════════════
# SLIDE 2 — TWO PATHS COMPARISON  (body ≥10pt)
# ═══════════════════════════════════════════════════════════

def build_s2():
    s = prs.slides.add_slide(blank_layout)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE
    add_header(s)
    add_title(s, 'AI4S 选题的两条路径：方法驱动 vs 问题驱动')
    add_assertion(s, '选题的第一步不是想 idea，而是想清楚你的起点——这决定了后续所有技术决策的方向')

    pw = Inches(5.4)                 # panel width
    py = Inches(1.88)                # panel top
    p1x = LM
    p2x = LM + pw + Inches(0.4)

    # ── PANEL A ──
    _rect(s, p1x, py, pw, Inches(3.9), fill=WHITE, border=USTC_BLUE, bw=Pt(2.5), radius=Inches(0.06))
    _rect(s, p1x, py, pw, Inches(0.45), fill=USTC_BLUE)
    _tb(s, p1x, py + Inches(0.05), pw, Inches(0.35),
        '路径 A：AI 方法驱动  |  Benchmark-Oriented', 14, WHITE, True, PP_ALIGN.CENTER)

    pa = [
        ('起点',     '我有一个更好的模型 / 方法', USTC_BLUE),
        ('核心逻辑', '现有方法在 benchmark 上存在局限\n→ 提出改进 → 超越 SOTA', DARK_GRAY),
        ('技术侧重', '模型架构创新、训练策略、通用性', DARK_GRAY),
        ('典型场景', 'QM9 / MD17 / OC20 刷榜\n新材料预测、通用 force field', DARK_GRAY),
        ('Gap 来源', 'AI 方法自身的局限性\n（如 GNN 的表达能力边界）', DARK_GRAY),
        ('核心风险', '刷榜与真科学问题脱节\n→ "手里有锤子，看什么都是钉子"', ACCENT_RED),
    ]
    _panel_rows(s, p1x + Inches(0.15), py + Inches(0.58), pw - Inches(0.3), pa,
                row_h=Inches(0.54), label_sz=10, value_sz=10)

    # ── PANEL B ──
    _rect(s, p2x, py, pw, Inches(3.9), fill=WHITE, border=GREEN, bw=Pt(2.5), radius=Inches(0.06))
    _rect(s, p2x, py, pw, Inches(0.45), fill=GREEN)
    _tb(s, p2x, py + Inches(0.05), pw, Inches(0.35),
        '路径 B：Science 问题驱动  |  Domain-Oriented', 14, WHITE, True, PP_ALIGN.CENTER)

    pb = [
        ('起点',     '有一个亟待解决的科学问题', GREEN),
        ('核心逻辑', '领域专家有一个认知空白\n→ 用 AI 填补空白 → 科学发现', DARK_GRAY),
        ('技术侧重', '物理约束、领域知识、可解释性', DARK_GRAY),
        ('典型场景', '帮材料学家发现新催化剂\n预测未知结构、解释实验现象', DARK_GRAY),
        ('Gap 来源', '科学认知的空白\n（某种材料的行为未被理解）', DARK_GRAY),
        ('核心风险', '技术贡献不够新颖\n→ 沦为"旧方法在新领域的应用"', ACCENT_RED),
    ]
    _panel_rows(s, p2x + Inches(0.15), py + Inches(0.58), pw - Inches(0.3), pb,
                row_h=Inches(0.54), label_sz=10, value_sz=10)

    # ── Bottom strip ──
    _strip(s, LM, Inches(5.98), Inches(11.6), Inches(0.65),
           '核心认知：两条路径不是互斥的。最优秀的研究在交集处——用 AI 方法创新解决真实科学问题。'
           '关键是你清醒地知道自己的起点在哪，选择承担哪类风险。',
           bg=LIGHT_BLUE, tc=USTC_BLUE, sz=11)

    add_footer(s, 2)


def _panel_rows(slide, x, y, w, items, row_h=Inches(0.54), label_sz=10, value_sz=10):
    for i, (label, value, clr) in enumerate(items):
        ry = y + i * row_h
        _tb(slide, x, ry, Inches(0.7), Inches(0.2), label, label_sz, clr, True)
        _tb(slide, x + Inches(0.73), ry, w - Inches(0.73), row_h + Inches(0.08),
            value, value_sz, DARK_GRAY, False, ls=1.35)


# ═══════════════════════════════════════════════════════════
# SLIDE 3 — GAP IDENTIFICATION  (body ≥10pt)
# ═══════════════════════════════════════════════════════════

def build_s3():
    s = prs.slides.add_slide(blank_layout)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE
    add_header(s)
    add_title(s, '不同路径，不同的 Gap 识别逻辑')
    add_assertion(s, '路径 A 找"方法的边界"——现有技术做不到什么；路径 B 找"科学的空白"——领域专家还不知道什么')

    top_y = Inches(1.88)
    half_w = Inches(5.5)
    gap_x = Inches(0.3)

    # ── LEFT: Path A ──
    _tb(s, LM, top_y, Inches(4.5), Inches(0.25),
        '▎路径 A：方法边界分析（找技术缺口）', 13, USTC_BLUE, True)
    a_methods = [
        ('技术图谱法', '画出领域内全部方法的能力矩阵\n→ 未被覆盖的区域 = 潜在 Gap', USTC_BLUE),
        ('SOTA 瓶颈法', '分析 benchmark 上方法的系统性失败模式\n→ 集体搞不定的 case = Gap', USTC_BLUE),
        ('理论边界法', '分析方法理论的表达能力上限\n（如 GNN WL-test 边界）→ Gap', USTC_BLUE),
    ]
    for i, (mt, md, mc) in enumerate(a_methods):
        mx = LM + Inches(0.08)
        my = top_y + Inches(0.38) + i * Inches(0.72)
        _rect(s, mx, my, Inches(0.05), Inches(0.62), fill=mc)
        _tb(s, mx + Inches(0.13), my, Inches(1.2), Inches(0.2), mt, 11, mc, True)
        _tb(s, mx + Inches(1.35), my, Inches(4.0), Inches(0.65), md, 10, DARK_GRAY, False, ls=1.3)

    # ── RIGHT: Path B ──
    r2x = LM + half_w + gap_x
    _tb(s, r2x, top_y, Inches(4.5), Inches(0.25),
        '▎路径 B：科学空白分析（找认知缺口）', 13, GREEN, True)
    b_methods = [
        ('领域需求访谈', '与领域专家深度对话\n→ 专家回答不了的 = 真问题', GREEN),
        ('文献对立面', '找两篇顶刊的结论冲突\n（实验 vs 计算的矛盾）→ Gap', GREEN),
        ('交叉组合法', 'A 领域的方法 × B 领域的问题\n→ 没人交叉过的 = Gap', GREEN),
    ]
    for i, (mt, md, mc) in enumerate(b_methods):
        mx = r2x + Inches(0.08)
        my = top_y + Inches(0.38) + i * Inches(0.72)
        _rect(s, mx, my, Inches(0.05), Inches(0.62), fill=mc)
        _tb(s, mx + Inches(0.13), my, Inches(1.2), Inches(0.2), mt, 11, mc, True)
        _tb(s, mx + Inches(1.35), my, Inches(4.0), Inches(0.65), md, 10, DARK_GRAY, False, ls=1.3)

    # ── BOTTOM: Gap test table ──
    table_top = Inches(4.28)
    _tb(s, LM, table_top - Inches(0.28), Inches(6), Inches(0.22),
        '▎两条路径通用的 Gap 检验标准', 12, ACCENT_RED, True)

    cw = [Inches(1.5), Inches(3.5), Inches(3.5), Inches(3.0)]
    tdata = [
        ['检验维度', '假 Gap', '真 Gap', '判断方法'],
        ['可替代性', '"还没人用 X 做过"\n→ 换个工具也能做',
         '"A 和 B 互补性未被利用"\n→ 互补是本质的',
         '不用 AI，这 Gap 还在吗？'],
        ['可证伪性', '实验必定成功\n没有失败空间',
         '任何实验结果都有意义\n（包括推翻假设）',
         '结果负面，有价值吗？'],
        ['问题来源', '从方法出发找问题\n"我有个锤子，找钉子"',
         '从问题出发找方法\n"有钉子，找合适的锤子"',
         '不用 AI，问题还存在吗？'],
    ]
    _table(s, LM, table_top, cw, tdata, font_sz=9.5)

    add_footer(s, 3)


# ═══════════════════════════════════════════════════════════
# SLIDE 4 — DECISION + CHECKLIST + MINI CASE  (body ≥10pt)
# ═══════════════════════════════════════════════════════════

def build_s4():
    s = prs.slides.add_slide(blank_layout)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE
    add_header(s)
    add_title(s, '选题决策：选择你的路径')
    add_assertion(s, '没有"正确"的路径——只有与你当前资源、阶段、目标匹配的路径')

    # ── LEFT: Decision Matrix ──
    left_x = LM
    left_w = Inches(6.5)
    dm_y = Inches(1.88)

    _tb(s, left_x, dm_y, Inches(4), Inches(0.22),
        '▎路径选择决策矩阵', 13, USTC_BLUE, True)

    dc = [Inches(1.8), Inches(2.3), Inches(2.3)]
    ddata = [
        ['评估维度', '适合路径 A (AI驱动)', '适合路径 B (科学驱动)'],
        ['团队构成', 'AI 团队，可独立完成', '有领域专家合作者'],
        ['计算资源', '充足，可跑大规模实验', '中等，重在分析验证'],
        ['发表目标', 'AI 顶会 (NeurIPS/ICML)', 'Science/Nature + AI 双栖'],
        ['研究阶段', '早期 career，建立声誉', '成熟期，追求领域 impact'],
        ['个人兴趣', '享受方法论创新', '享受解决实际问题'],
        ['风险偏好', '接受 benchmark 风险', '接受"不够 novel"风险'],
    ]
    _table(s, left_x, dm_y + Inches(0.32), dc, ddata, font_sz=9.5)

    # ── RIGHT: Self-check ──
    right_x = left_x + left_w + Inches(0.3)
    right_w = Inches(5.2)
    ck_y = Inches(1.88)

    _tb(s, right_x, ck_y, right_w, Inches(0.22),
        '▎选题前五问（自查清单）', 13, ORANGE, True)

    qs = [
        ('① 我的 Gap 在哪个空间？\n   方法边界 → 路径A    科学空白 → 路径B', DARK_GRAY),
        ('② 如果不用 AI，这个问题还存在吗？\n   不 → 假问题    是 → 真科学问题', ACCENT_RED),
        ('③ 我能画出"能力边界图"吗？\n   列出已知方法 × 已知问题的矩阵，找空白格', DARK_GRAY),
        ('④ 我的实验设计能接受失败吗？\n   如果结果负面，它是否有分析价值？', ACCENT_RED),
        ('⑤ 我缺什么？（合作者 / 算力 / 领域知识）\n   缺少的关键资源影响路径可行性', DARK_GRAY),
    ]
    for i, (q, clr) in enumerate(qs):
        qy = ck_y + Inches(0.35) + i * Inches(0.68)
        _tb(s, right_x + Inches(0.08), qy, right_w - Inches(0.15), Inches(0.62),
            q, 10, clr, (i == 1 or i == 3), ls=1.4)

    # ── BOTTOM: Mini Case ──
    case_y = Inches(5.65)
    _line(s, LM, case_y, Inches(11.6), LIGHT_GRAY, Pt(1.5))
    _tb(s, LM, case_y + Inches(0.04), Inches(4), Inches(0.2),
        '▎Mini Case：CrysLLMGen 的路径选择', 11, USTC_BLUE, True)

    _tb(s, LM, case_y + Inches(0.26), Inches(11.6), Inches(0.75),
        '起点方法驱动（路径A）：发现 LLM 和 Diffusion 各自有明确的"做不到"（LLM 缺精确坐标，Diffusion 缺化学先验），'
        '两者的互补性构成了一个真 Gap。随后融入问题驱动（路径B）：晶体材料生成是材料科学家的真实需求。'
        '最终定位在两条路径的交集——用混合框架这一架构创新解决真实科学问题。NeurIPS 2025。',
        10, DARK_GRAY, False, ls=1.35)

    add_footer(s, 4)


# ═══════════════════════════════════════════════════════════

build_s1(); build_s2(); build_s3(); build_s4()
out = '/zhdd/home/hengzhang/code/crysllmgen-main/AI4S科研选题方法论_张恒.pptx'
prs.save(out)
print(f'OK → {out}  ({len(prs.slides)} slides)')
